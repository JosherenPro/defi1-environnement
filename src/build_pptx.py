#!/usr/bin/env python3
"""
Génère le rapport PowerPoint (<=10 pages, FR) du Défi 1 Environnement.
Design & Identité : Togo AI Lab / Service Public Togo (Vert #0B4F4A, Or #F4B400, Poppins).
"""
import os, json
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import geopandas as gpd
from matplotlib.lines import Line2D
from matplotlib.patches import Patch
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
AST  = os.path.join(BASE, "dashboard", "assets")
OUT  = os.path.join(BASE, "rapport.pptx")
IMG  = os.path.join(BASE, "reports", "assets")
os.makedirs(IMG, exist_ok=True)
LOGO_IMG = os.path.join(IMG, "logo-datalab.png")

data = json.load(open(os.path.join(AST, "data.json"), encoding="utf-8"))
R, F, W, C, S = data["regions"], data["fri_classes"], data["water_sales"], data["coso"], data["summary"]
P = data.get("priority_cantons", {})
CONTROL_COUNTS = C.get("control_class_counts", {})
CONTROL_PRIORITY = int(S.get("coso_control_priority_count", CONTROL_COUNTS.get("Priorité de contrôle", 0)))
COSO_MISSING_COORD_PCT = round(100 - float(S["coso_coord_quality_pct"]), 1)
MAINTENANCE_MISSING_TOTAL = int(sum(R.get("maintenance_missing_coso", [])))

def fr1(value):
    return f"{float(value):.1f}".replace(".", ",")

TEAL   = "#0B4F4A"
GOLD   = "#F4B400"
TURQ   = "#14877D"
ORG    = "#D9622B"
PURP   = "#7A4FA0"
INK    = "#16302C"
MUTED  = "#5D726E"
BG_LIGHT = "#F4F7F6"
CARD_BG = "#FFFFFF"
BORDER_CLR = "#D8E3E1"

REGION_COLORS = {"Maritime": TEAL, "Plateaux": TURQ, "Centrale": GOLD, "Kara": ORG, "Savanes": PURP}
FONT_TEXT = "Poppins"
FONT_TITLE = "Poppins"

plt.rcParams.update({
    "font.family": "sans-serif",
    "font.sans-serif": ["Poppins", "DejaVu Sans", "Liberation Sans"],
    "axes.edgecolor": "#D8E3E1",
    "axes.linewidth": 0.8,
})

def save_bars(x, y, colors, fname, title, ylabel=""):
    fig, ax = plt.subplots(figsize=(6.4, 3.4), dpi=180)
    bars = ax.bar(x, y, color=colors, width=0.55, edgecolor="none", zorder=3)
    for bar_item in bars:
        height = bar_item.get_height()
        val_str = f"{height:g}" if isinstance(height, (int, float)) and height == int(height) else f"{height:.1f}"
        ax.annotate(val_str,
                    xy=(bar_item.get_x() + bar_item.get_width() / 2, height),
                    xytext=(0, 4), textcoords="offset points",
                    ha='center', va='bottom', fontsize=9, fontweight='bold', color=INK)
    ax.set_title(title, fontsize=12, color=TEAL, fontweight="bold", pad=12, loc="left")
    if ylabel:
        ax.set_ylabel(ylabel, fontsize=9, color=MUTED)
    ax.tick_params(axis="both", which="major", labelsize=9, colors=INK)
    ax.spines[["top", "right"]].set_visible(False)
    ax.grid(axis="y", linestyle="--", alpha=0.4, color="#cccccc", zorder=0)
    plt.xticks(rotation=15, ha="right")
    plt.tight_layout()
    p = os.path.join(IMG, fname)
    fig.savefig(p, bbox_inches="tight", dpi=180)
    plt.close(fig)
    return p

def save_grouped_bars(x, y1, y2, l1, l2, c1, c2, fname, title):
    fig, ax = plt.subplots(figsize=(6.4, 3.4), dpi=180)
    indices = np.arange(len(x))
    width = 0.35
    b1 = ax.bar(indices - width/2, y1, width, label=l1, color=c1, zorder=3)
    b2 = ax.bar(indices + width/2, y2, width, label=l2, color=c2, zorder=3)
    for bar_item in b1:
        h = bar_item.get_height()
        if h > 0:
            ax.annotate(f"{int(h)}", xy=(bar_item.get_x() + bar_item.get_width()/2, h),
                        xytext=(0, 3), textcoords="offset points", ha='center', va='bottom', fontsize=8, color=INK)
    for bar_item in b2:
        h = bar_item.get_height()
        if h > 0:
            ax.annotate(f"{int(h)}", xy=(bar_item.get_x() + bar_item.get_width()/2, h),
                        xytext=(0, 3), textcoords="offset points", ha='center', va='bottom', fontsize=8, color=INK)
    ax.set_title(title, fontsize=12, color=TEAL, fontweight="bold", pad=12, loc="left")
    ax.set_xticks(indices)
    ax.set_xticklabels(x, fontsize=9, color=INK)
    ax.tick_params(axis="y", labelsize=8, colors=INK)
    ax.legend(frameon=True, facecolor="white", edgecolor="#D8E3E1", fontsize=8)
    ax.spines[["top", "right"]].set_visible(False)
    ax.grid(axis="y", linestyle="--", alpha=0.4, color="#cccccc", zorder=0)
    plt.tight_layout()
    p = os.path.join(IMG, fname)
    fig.savefig(p, bbox_inches="tight", dpi=180)
    plt.close(fig)
    return p

def save_pie(labels, vals, colors, fname, title):
    fig, ax = plt.subplots(figsize=(4.8, 3.4), dpi=180)
    wedges, texts, autotexts = ax.pie(
        vals, labels=labels, colors=colors, autopct="%1.0f%%", startangle=90,
        pctdistance=0.7, wedgeprops=dict(width=0.45, edgecolor="white", linewidth=2)
    )
    for t in texts:
        t.set_fontsize(9)
        t.set_color(INK)
    for at in autotexts:
        at.set_fontsize(9)
        at.set_fontweight("bold")
        at.set_color("white")
    ax.set_title(title, fontsize=11, color=TEAL, fontweight="bold", pad=12, loc="center")
    plt.tight_layout()
    p = os.path.join(IMG, fname)
    fig.savefig(p, bbox_inches="tight", dpi=180)
    plt.close(fig)
    return p

def save_line(years, series, fname, title):
    fig, ax = plt.subplots(figsize=(6.8, 3.4), dpi=180)
    colors = [TEAL, TURQ, GOLD, ORG, PURP, "#2E7D32", "#D32F2F", "#0288D1", "#78909C"]
    for i, (cat, vals) in enumerate(series):
        ax.plot(years, vals, marker="o", markersize=4, linewidth=1.8, label=cat, color=colors[i % len(colors)])
    ax.set_title(title, fontsize=12, color=TEAL, fontweight="bold", pad=12, loc="left")
    ax.legend(fontsize=7, ncol=2, loc="upper left", frameon=True, facecolor="white", edgecolor="#D8E3E1")
    ax.tick_params(labelsize=8, colors=INK)
    ax.spines[["top", "right"]].set_visible(False)
    ax.grid(True, linestyle="--", alpha=0.3, color="#cccccc")
    plt.tight_layout()
    p = os.path.join(IMG, fname)
    fig.savefig(p, bbox_inches="tight", dpi=180)
    plt.close(fig)
    return p

def save_map(fname):
    cantons = gpd.read_file(os.path.join(AST, "cantons.geojson"))
    points = gpd.read_file(os.path.join(AST, "points.geojson"))
    fri_colors = {"Faible": TEAL, "Moyen": GOLD, "Élevé": ORG}
    fig, ax = plt.subplots(figsize=(6.4, 4.55), dpi=180)
    for fri_class, color in fri_colors.items():
        subset = cantons[cantons["fri_class"].astype(str) == fri_class]
        if not subset.empty:
            subset.plot(ax=ax, color=color, alpha=0.52, edgecolor="white", linewidth=0.22)
    for source, color in [("TdE", TEAL), ("COSO", PURP)]:
        subset = points[points["src"].astype(str) == source]
        if not subset.empty:
            subset.plot(ax=ax, color=color, markersize=9, alpha=0.9,
                        edgecolor="white", linewidth=0.35)
    ax.set_axis_off()
    ax.set_title("Ouvrages géolocalisés et exposition FRI", fontsize=12,
                 color=TEAL, fontweight="bold", pad=8, loc="left")
    legend_handles = [
        Patch(facecolor=TEAL, alpha=0.52, edgecolor="none", label="FRI faible"),
        Patch(facecolor=GOLD, alpha=0.52, edgecolor="none", label="FRI moyen"),
        Patch(facecolor=ORG, alpha=0.52, edgecolor="none", label="FRI élevé"),
        Line2D([0], [0], marker="o", color="w", markerfacecolor=TEAL,
               markeredgecolor="white", markersize=6, label="TdE"),
        Line2D([0], [0], marker="o", color="w", markerfacecolor=PURP,
               markeredgecolor="white", markersize=6, label="COSO"),
    ]
    fig.legend(handles=legend_handles, loc="center right",
               bbox_to_anchor=(0.98, 0.50), ncol=1, fontsize=7,
               frameon=True, facecolor="white", edgecolor="#D8E3E1",
               framealpha=0.95)
    plt.tight_layout()
    p = os.path.join(IMG, fname)
    fig.savefig(p, dpi=180, facecolor="white")
    plt.close(fig)
    return p

# Graphiques
fig_pts = save_grouped_bars(
    R["labels"], R["n_tde"], R["n_coso"],
    "TdE (Châteaux d'eau/Forages)", "COSO (Microprojets Nord)",
    TEAL, ORG, "pts_grouped.png", "Répartition des ouvrages par source & région"
)
fig_dens = save_bars(R["labels"], R["points_per_100k"], [REGION_COLORS[r] for r in R["labels"]], "dens.png", "Densité d'équipements (points / 100 000 hab.)", "Points / 100k hab.")
maint_labels = [r for r, n in zip(R["labels"], R["n_coso"]) if n > 0]
maint_values = [v for v, n in zip(R["maint_rate_coso"], R["n_coso"]) if n > 0]
fig_maint = save_bars(maint_labels, maint_values, [REGION_COLORS[r] for r in maint_labels], "maint.png", "Projets COSO avec plan d'entretien (%)", "% avec plan")
fig_pie = save_pie(F["class"], F["n_cantons"], [TEAL, GOLD, ORG], "fripie.png", "Répartition des cantons par classe FRI")
years = W["years"]
series = [(W["categories"][i], W["matrix"][i]) for i in range(len(W["categories"]))]
fig_sales = save_line(years, series, "sales.png", "Évolution des ventes d'eau facturées TdE (m³)")
fig_map = save_map("map_points_fri.png")

COSO_MISSING_COORD_N = int(S["n_coso"] - S["coso_with_coord"])

# PowerPoint
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def rgb(h):
    return RGBColor.from_string(h.lstrip("#"))

def remove_shape_shadow(shape):
    style = shape._element.find(qn("p:style"))
    if style is not None:
        effect_ref = style.find(qn("a:effectRef"))
        if effect_ref is not None:
            style.remove(effect_ref)

TEALC  = rgb(TEAL)
GOLDC  = rgb(GOLD)
TURQC  = rgb(TURQ)
ORGC   = rgb(ORG)
PURPC  = rgb(PURP)
INKC   = rgb(INK)
MUTEDC = rgb(MUTED)
WHITE  = rgb("#FFFFFF")
CARDC  = rgb(CARD_BG)
BORDC  = rgb(BORDER_CLR)
LIGHTC = rgb(BG_LIGHT)

def add_rect(slide, x, y, w, h, fill_color, border_color=None, border_width=1):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill_color
    if border_color:
        sp.line.color.rgb = border_color
        sp.line.width = Pt(min(border_width, 0.6))
    else:
        sp.line.fill.background()
    remove_shape_shadow(sp)
    sp.shadow.inherit = False
    return sp

def add_card(slide, x, y, w, h, bg_color=WHITE, border_color=BORDC, corner_radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shape_type, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = bg_color
    if border_color:
        sp.line.color.rgb = border_color
        sp.line.width = Pt(0.6)
    else:
        sp.line.fill.background()
    remove_shape_shadow(sp)
    sp.shadow.inherit = False
    return sp

def txt(slide, x, y, w, h, text, size=13, color=INKC, bold=False, align=PP_ALIGN.LEFT, font=FONT_TEXT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(text)
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    f.name = font
    return tb

def datalab_logo(slide, x, y, w, h):
    if os.path.exists(LOGO_IMG):
        slide.shapes.add_picture(LOGO_IMG, x, y, width=w, height=h)

def header(slide, category_tag, title_text):
    add_rect(slide, 0, 0, SW, Inches(1.15), TEALC)
    add_rect(slide, 0, Inches(1.15), SW, Pt(4), GOLDC)
    datalab_logo(slide, Inches(12.55), Inches(0.18), Inches(0.32), Inches(0.50))
    txt(slide, Inches(0.5), Inches(0.12), Inches(11.8), Inches(0.28), category_tag.upper(), 10, GOLDC, bold=True, font=FONT_TITLE)
    txt(slide, Inches(0.5), Inches(0.38), Inches(11.8), Inches(0.70), title_text, 22, WHITE, bold=True, font=FONT_TITLE)

def footer(slide, slide_num):
    add_rect(slide, Inches(0.5), Inches(7.0), Inches(12.333), Pt(0.7), BORDC)
    txt(slide, Inches(0.5), Inches(7.08), Inches(10), Inches(0.32),
        "TOGO AI LAB · Data Challenge Environnement (Défi 01) | Données Ouvertes TdE · COSO · INSEED · ISRI-TG", 9, MUTEDC, font=FONT_TEXT)
    txt(slide, Inches(12.333), Inches(7.08), Inches(0.5), Inches(0.32), str(slide_num), 9, MUTEDC, align=PP_ALIGN.RIGHT, bold=True, font=FONT_TEXT)

# Slide 1
s1 = prs.slides.add_slide(BLANK)
add_rect(s1, 0, 0, SW, SH, TEALC)
add_rect(s1, 0, Inches(4.5), SW, Pt(5), GOLDC)
datalab_logo(s1, Inches(11.7), Inches(0.6), Inches(0.8), Inches(1.25))
txt(s1, Inches(0.8), Inches(1.2), Inches(10.5), Inches(0.4), "TOGO AI LAB · DATA CHALLENGE ENVIRONNEMENT (DÉFI 01)", 13, GOLDC, bold=True)
txt(s1, Inches(0.8), Inches(1.7), Inches(11.2), Inches(1.5), "Accès à l'Eau Potable au Togo", 42, WHITE, bold=True)
txt(s1, Inches(0.8), Inches(3.2), Inches(11.0), Inches(1.1),
    "Diagnostic spatial des infrastructures, évaluation de la durabilité du parc, analyse de la pression démographique et cartographie du risque d'inondation (FRI).", 16, rgb("#E0EFEA"))
chip_data = [
    ("285 Ouvrages", "TdE (67) + COSO (218)"),
    ("388 Cantons", "Évalués FRI (ISRI-TG)"),
    ("Livraison Complète", "Dashboard Streamlit + PPTX")
]
cx = Inches(0.8)
for title_c, desc_c in chip_data:
    add_card(s1, cx, Inches(4.85), Inches(3.6), Inches(1.2), bg_color=rgb("#083E3A"), border_color=GOLDC)
    txt(s1, cx + Inches(0.2), Inches(5.0), Inches(3.2), Inches(0.4), title_c, 15, GOLDC, bold=True)
    txt(s1, cx + Inches(0.2), Inches(5.4), Inches(3.2), Inches(0.5), desc_c, 11, WHITE)
    cx += Inches(3.85)
txt(s1, Inches(0.8), Inches(6.65), Inches(11.7), Inches(0.4),
    "Sources de données ouvertes : TdE (Inventaire & Ventes) · Projet COSO (Microprojets) · INSEED/RGPH 2010 · ISRI-TG (Indices FSI/FRI)", 10, rgb("#A0C4C0"))

# Slide 2
s2 = prs.slides.add_slide(BLANK)
header(s2, "CADRE GÉNÉRAL", "Contexte & Démarche Analytique")
add_card(s2, Inches(0.5), Inches(1.4), Inches(6.5), Inches(5.4), bg_color=WHITE)
txt(s2, Inches(0.7), Inches(1.6), Inches(6.1), Inches(0.4), "Objectifs & Méthodologie du Défi", 15, TEALC, bold=True)
context_items = [
    ("Diagnostic Intégré", "Évaluer l'accès à l'eau potable au Togo à travers le croisement de 5 jeux de données ouverts nationaux."),
    ("Pipeline ETL Python", "Normalisation géographique sous pandas/geopandas, jointure spatiale points × cantons FRI et calcul d'indicateurs régionaux."),
    ("Proxy de Durabilité", "En l'absence de statut opérationnel publié, utilisation du taux de plan d'entretien documenté (COSO) et d'un score de priorité de contrôle."),
    ("Livrables Réglementaires", "Déploiement d'un dashboard analytique interactif (Streamlit) et présent rapport d'analyse synthétique (≤ 10 pages).")
]
cy = Inches(2.1)
for title_i, desc_i in context_items:
    add_rect(s2, Inches(0.7), cy, Pt(4), Inches(0.8), TURQC)
    txt(s2, Inches(0.85), cy, Inches(5.9), Inches(0.3), title_i, 13, TEALC, bold=True)
    txt(s2, Inches(0.85), cy + Inches(0.28), Inches(5.9), Inches(0.65), desc_i, 11, INKC)
    cy += Inches(1.15)

kx = Inches(7.3)
add_card(s2, kx, Inches(1.4), Inches(5.5), Inches(5.4), bg_color=LIGHTC)
txt(s2, kx + Inches(0.3), Inches(1.6), Inches(4.9), Inches(0.4), "Chiffres Clés du Diagnostic", 15, TEALC, bold=True)
kpi_list = [
    ("285", "Ouvrages hydrauliques recensés au total", TEAL),
    ("150", f"Ouvrages avec géolocalisation exploitable ({fr1(S['n_points_geoloc_ok'] / S['n_points_total'] * 100)} %)", TURQ),
    (fr1(S['coso_coord_quality_pct']) + " %", f"Points COSO avec coordonnées valides ({COSO_MISSING_COORD_N} à corriger)", ORG),
    (fr1(S['coso_maint_overall_pct']) + " %", "Projets COSO avec plan d'entretien documenté", ORG),
    (f"{CONTROL_PRIORITY}", "Ouvrages COSO classés en priorité de contrôle", ORG)
]
ky = Inches(2.15)
for val_k, lbl_k, col_k in kpi_list:
    add_card(s2, kx + Inches(0.3), ky, Inches(4.9), Inches(0.85), bg_color=WHITE, border_color=rgb("#E0EBE8"))
    txt(s2, kx + Inches(0.5), ky + Inches(0.12), Inches(1.4), Inches(0.6), val_k, 22, rgb(col_k), bold=True)
    txt(s2, kx + Inches(1.9), ky + Inches(0.15), Inches(3.2), Inches(0.6), lbl_k, 11, INKC)
    ky += Inches(0.95)
footer(s2, 2)

# Slide 3
s3 = prs.slides.add_slide(BLANK)
header(s3, "OBJECTIF 1", "Cartographie & Répartition Spatiale des Ouvrages")
s3.shapes.add_picture(fig_map, Inches(0.5), Inches(1.4), width=Inches(6.4))
add_card(s3, Inches(7.1), Inches(1.4), Inches(5.7), Inches(5.4), bg_color=WHITE)
txt(s3, Inches(7.3), Inches(1.6), Inches(5.3), Inches(0.4), "Constats & Qualité des Données Spatiales", 15, TEALC, bold=True)
bullet_s3 = [
    ("Hétérogénéité territoriale extrême",
     f"L'inventaire TdE est quasi-exclusivement centré sur la région Maritime ({R['n_tde'][0]}/{S['n_tde']} ouvrages, soit 97 %). Les microprojets COSO couvrent le Nord (Savanes: 192, Kara: 13, Centrale: 13)."),
    ("Défi majeur de géolocalisation",
     f"Seuls {S['n_points_geoloc_ok']} des {S['n_points_total']} ouvrages disposent de coordonnées valides. {fr1(COSO_MISSING_COORD_PCT)} % des points COSO ({COSO_MISSING_COORD_N}/{S['n_coso']}) présentent des coordonnées nulles (0,0) et ne peuvent pas être cartographiés."),
    ("Conséquence pour le pilotage",
     "Les cartes actuelles reflètent la localisation des inventaires disponibles et non l'exhaustivité de la couverture nationale en eau potable.")
]
by3 = Inches(2.2)
for b_title, b_desc in bullet_s3:
    add_rect(s3, Inches(7.3), by3, Pt(4), Inches(0.9), ORGC if "Défi" in b_title else TURQC)
    txt(s3, Inches(7.45), by3, Inches(5.1), Inches(0.3), b_title, 12, TEALC, bold=True)
    txt(s3, Inches(7.45), by3 + Inches(0.28), Inches(5.1), Inches(0.9), b_desc, 11, INKC)
    by3 += Inches(1.5)
footer(s3, 3)

# Slide 4
s4 = prs.slides.add_slide(BLANK)
header(s4, "OBJECTIF 2", "État du Parc & Proxy de Maintenance")
s4.shapes.add_picture(fig_maint, Inches(0.5), Inches(1.4), width=Inches(6.4))
add_card(s4, Inches(7.1), Inches(1.4), Inches(5.7), Inches(5.4), bg_color=WHITE)
txt(s4, Inches(7.3), Inches(1.6), Inches(5.3), Inches(0.4), "Diagnostic de Durabilité & Score de Contrôle", 15, TEALC, bold=True)
maint_bullets = [
    ("Absence de statut opérationnel publié",
     "Les données ouvertes TdE et COSO ne contiennent pas les champs 'fonctionnel', 'en panne' ou 'abandonné'. Le taux de fonctionnalité réel ne peut donc pas être calculé directement."),
    ("Taux de plan d'entretien très faible",
     f"Seuls {fr1(S['coso_maint_overall_pct'])} % des projets COSO disposent d'un plan d'entretien documenté (100 % en Kara, {fr1(R['maint_rate_coso'][4])} % en Centrale, {fr1(R['maint_rate_coso'][2])} % en Savanes)."),
    ("Score indicatif de priorité de contrôle",
     f"Un score composite (0 à 100) a été modélisé : absence de plan (45%), statut non définitif (25%), absence de remise (15%) et FRI (15%). Il classe {CONTROL_PRIORITY} ouvrages en 'Priorité de contrôle' pour cibler les audits terrain.")
]
m_by = Inches(2.2)
for b_title, b_desc in maint_bullets:
    add_rect(s4, Inches(7.3), m_by, Pt(4), Inches(0.95), ORGC if "Score" in b_title or "Absence" in b_title else TURQC)
    txt(s4, Inches(7.45), m_by, Inches(5.1), Inches(0.3), b_title, 12, TEALC, bold=True)
    txt(s4, Inches(7.45), m_by + Inches(0.28), Inches(5.1), Inches(0.9), b_desc, 11, INKC)
    m_by += Inches(1.5)
footer(s4, 4)

# Slide 5
s5 = prs.slides.add_slide(BLANK)
header(s5, "OBJECTIF 3", "Pression Démographique vs Infrastructures")
s5.shapes.add_picture(fig_dens, Inches(0.5), Inches(1.4), width=Inches(6.4))
add_card(s5, Inches(7.1), Inches(1.4), Inches(5.7), Inches(5.4), bg_color=WHITE)
txt(s5, Inches(7.3), Inches(1.6), Inches(5.3), Inches(0.4), "Analyse de la Pression & Équité Territoriale", 15, TEALC, bold=True)
demo_bullets = [
    ("Forte disparité de densité apparente",
     f"La densité varie de {R['points_per_100k'][2]:.2f} points / 100k hab. en Savanes (forte concentration COSO) à seulement {R['points_per_100k'][1]:.2f} en Plateaux et {R['points_per_100k'][0]:.2f} en Maritime (selon inventaires publiés)."),
    ("Base démographique INSEED 2010",
     "L'analyse s'appuie sur la population légale du RGPH 2010. Le secteur des Plateaux affiche le déficit relatif d'équipement documenté le plus prononcé."),
    ("Vigilance d'interprétation",
     "La densité d'ouvrages ne mesure pas le taux de desserte effective : la capacité, le débit et le rayon de couverture de chaque point d'eau ne sont pas renseignés dans les jeux ouverts.")
]
d_by = Inches(2.2)
for b_title, b_desc in demo_bullets:
    add_rect(s5, Inches(7.3), d_by, Pt(4), Inches(0.95), TURQC)
    txt(s5, Inches(7.45), d_by, Inches(5.1), Inches(0.3), b_title, 12, TEALC, bold=True)
    txt(s5, Inches(7.45), d_by + Inches(0.28), Inches(5.1), Inches(0.9), b_desc, 11, INKC)
    d_by += Inches(1.5)
footer(s5, 5)

# Slide 6
s6 = prs.slides.add_slide(BLANK)
header(s6, "OBJECTIF 4", "Vulnérabilité & Risque d'Inondation (FRI)")
s6.shapes.add_picture(fig_pie, Inches(0.5), Inches(1.4), width=Inches(4.8))
add_card(s6, Inches(5.5), Inches(1.4), Inches(7.3), Inches(5.4), bg_color=WHITE)
txt(s6, Inches(5.7), Inches(1.6), Inches(6.9), Inches(0.4), "Croisement Spatiale : Cantons FRI × Ouvrages", 15, TEALC, bold=True)
fri_kpis = [
    ("75 Cantons", f"En risque 'Élevé' (FRI > 0.13), regroupant {S['population_high_FRI']:,} habitants.", ORG),
    ("77 Ouvrages", "Géolocalisés situés directement dans des cantons FRI 'Élevé' (47 en Maritime, 27 en Savanes).", ORG),
    ("Priorité Résilience", "Surélévation des têtes de forages, dalles étanches et protection contre la contamination des eaux de crue.", TEAL)
]
f_by = Inches(2.2)
for fk_val, fk_desc, fk_col in fri_kpis:
    add_card(s6, Inches(5.7), f_by, Inches(6.9), Inches(1.2), bg_color=LIGHTC, border_color=rgb("#E0EBE8"))
    txt(s6, Inches(5.9), f_by + Inches(0.15), Inches(2.2), Inches(0.4), fk_val, 16, rgb(fk_col), bold=True)
    txt(s6, Inches(5.9), f_by + Inches(0.55), Inches(6.5), Inches(0.6), fk_desc, 11, INKC)
    f_by += Inches(1.45)
footer(s6, 6)

# Slide 7
s7 = prs.slides.add_slide(BLANK)
header(s7, "ANALYSE COMPLÉMENTAIRE", "Structure & Évolution des Ventes d'Eau TdE")
s7.shapes.add_picture(fig_sales, Inches(0.5), Inches(1.4), width=Inches(6.8))
add_card(s7, Inches(7.5), Inches(1.4), Inches(5.3), Inches(5.4), bg_color=WHITE)
txt(s7, Inches(7.7), Inches(1.6), Inches(4.9), Inches(0.4), "Analyse des Ventes d'Eau (2018-2022)", 15, TEALC, bold=True)
sales_bullets = [
    ("Volume global 2022",
     f"Le volume total facturé par la TdE s'élève à {int(S['total_2022_m3']):,} m³ en 2022."),
    ("Prédominance de l'usage industriel",
     "Les catégories 'Forage Usage Autres' (141k m³) et 'Zone Franche' (126k m³) représentent la majorité des volumes d'eau distribués."),
    ("Consommation ménages marginale",
     "Les abonnements 'Concessions TdE' (63k m³) et 'Collectivités' (35k m³) ne représentent qu'une faible part du volume total, traduisant un besoin d'extension des branchements sociaux domestiques.")
]
s_by = Inches(2.2)
for b_title, b_desc in sales_bullets:
    add_rect(s7, Inches(7.7), s_by, Pt(4), Inches(0.95), TURQC)
    txt(s7, Inches(7.85), s_by, Inches(4.7), Inches(0.3), b_title, 12, TEALC, bold=True)
    txt(s7, Inches(7.85), s_by + Inches(0.28), Inches(4.7), Inches(0.9), b_desc, 11, INKC)
    s_by += Inches(1.5)
footer(s7, 7)

# Slide 8
s8 = prs.slides.add_slide(BLANK)
header(s8, "PLAN D'ACTION", "Recommandations opérationnelles — R1 à R7")
reco_cards = [
    ("R1", "MAINTENANCE IMMÉDIATE",
     f"Auditer les {MAINTENANCE_MISSING_TOTAL} ouvrages COSO sans plan documenté et programmer une visite trimestrielle avec statut panne/abandon.", TEAL),
    ("R2", "SUIVI DE FONCTIONNALITÉ",
     "Relever trimestriellement l'état de chaque ouvrage et publier les champs fonctionnel, en panne, abandonné, date et cause ; le score COSO cible les premières visites.", TURQ),
    ("R3", "QUALITÉ & OPEN DATA",
     f"Corriger les {COSO_MISSING_COORD_N} coordonnées nulles (0,0) et interconnecter les bases TdE/COSO sur opendata.gouv.tg avec un schéma géographique commun.", TURQ),
    ("R4", "NOUVEAUX FORAGES",
     "Cibler les cantons à score de priorité élevé, en croisant population, faible densité d'ouvrages et FRI élevé.", PURP),
    ("R5", "RÉSILIENCE",
     f"Surélever les équipements électriques, protéger les têtes de forage et prévoir drainage/accès de secours pour les {S['points_high_FRI']} ouvrages en FRI élevé.", ORG),
    ("R6", "ÉQUITÉ TERRITORIALE",
     "Compléter la couverture hors Maritime et hors corridor COSO des Savanes, après étude hydrogéologique et test de qualité de l'eau.", GOLD),
    ("R7", "ACCÈS DOMESTIQUE",
     "Flécher les investissements vers les branchements sociaux et points de desserte, les usages industriels/zone franche dominant les volumes TdE 2022.", TEAL),
]
for i, (r_code, r_title, r_desc, r_col) in enumerate(reco_cards):
    row, col = divmod(i, 2)
    cx = Inches(0.5) + col * Inches(6.25)
    cy = Inches(1.42) + row * Inches(1.34)
    add_card(s8, cx, cy, Inches(6.08), Inches(1.16), bg_color=WHITE, border_color=None)
    add_rect(s8, cx, cy, Inches(0.14), Inches(1.16), rgb(r_col))
    add_rect(s8, cx + Inches(0.28), cy + Inches(0.15), Inches(0.52), Inches(0.46), rgb(r_col))
    txt(s8, cx + Inches(0.29), cy + Inches(0.19), Inches(0.50), Inches(0.30), r_code, 11, WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s8, cx + Inches(1.0), cy + Inches(0.12), Inches(4.75), Inches(0.26), r_title, 11, TEALC, bold=True)
    txt(s8, cx + Inches(1.0), cy + Inches(0.43), Inches(4.75), Inches(0.61), r_desc, 9.5, INKC)
txt(s8, Inches(0.65), Inches(6.82), Inches(12.0), Inches(0.18),
    "Important : le score COSO sert au ciblage des visites ; il ne constitue pas un taux de panne.", 9, MUTEDC)
footer(s8, 8)

# Slide 9
s9 = prs.slides.add_slide(BLANK)
header(s9, "CONCLUSION", "Conclusion & Feuille de Route Gouvernance")
add_card(s9, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.4), bg_color=WHITE)
txt(s9, Inches(0.7), Inches(1.6), Inches(5.6), Inches(0.4), "Bilan Synthétique du Diagnostic", 15, TEALC, bold=True)
concl_items = [
    ("Couverture Hétérogène", "Forte disparité spatiale entre la région Maritime (TdE) et le corridor Nord (COSO)."),
    ("Durabilité à Sécuriser", "Seuls 20.6 % des ouvrages COSO ont un plan d'entretien documenté, nécessitant un suivi rigoureux."),
    ("Vulnérabilité Climatique", "77 ouvrages géolocalisés sont exposés à un risque d'inondation élevé (FRI > 0.13)."),
    ("Livrables Opérationnels", "Un dashboard interactif fonctionnel (.zip) et ce rapport d'analyse pour guider les décideurs.")
]
cy9 = Inches(2.15)
for c_t, c_d in concl_items:
    add_rect(s9, Inches(0.7), cy9, Pt(4), Inches(0.85), TEALC)
    txt(s9, Inches(0.85), cy9, Inches(5.4), Inches(0.3), c_t, 12, TEALC, bold=True)
    txt(s9, Inches(0.85), cy9 + Inches(0.28), Inches(5.4), Inches(0.6), c_d, 11, INKC)
    cy9 += Inches(1.15)

add_card(s9, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.4), bg_color=LIGHTC)
txt(s9, Inches(7.0), Inches(1.6), Inches(5.6), Inches(0.4), "Feuille de Route pour la Gouvernance", 15, TEALC, bold=True)
roadmap_items = [
    ("Étape 1 : Standardisation", "Instaurer un schéma géographique commun (GPS, codes cantons) pour tous les projets d'eau."),
    ("Étape 2 : Campagne Terrain", "Déployer un contrôle physique sur les 115 ouvrages identifiés en priorité de contrôle."),
    ("Étape 3 : Interconnexion Portail", "Mettre à jour et unifier les jeux TdE, COSO et Ministère sur l'écosystème opendata.gouv.tg."),
    ("Étape 4 : Suivi d'Impact", "Mesurer le taux de desserte effectif et la fonctionnalité réelle au niveau canton.")
]
ry9 = Inches(2.15)
for r_t, r_d in roadmap_items:
    add_card(s9, Inches(7.0), ry9, Inches(5.6), Inches(0.9), bg_color=WHITE, border_color=rgb("#E0EBE8"))
    txt(s9, Inches(7.2), ry9 + Inches(0.12), Inches(5.2), Inches(0.3), r_t, 12, TURQC, bold=True)
    txt(s9, Inches(7.2), ry9 + Inches(0.40), Inches(5.2), Inches(0.45), r_d, 11, INKC)
    ry9 += Inches(1.15)
footer(s9, 9)

prs.save(OUT)
print("Rapport PowerPoint généré :", OUT, "| slides :", len(prs.slides._sldIdLst))
